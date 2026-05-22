"""Generate 2 intro slides: RULER + Lost in the Middle."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --- Theme colors ---
BG_DARK = RGBColor(0x0B, 0x14, 0x26)
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_MUTED = RGBColor(0xB8, 0xC5, 0xD6)
ACCENT_CYAN = RGBColor(0x00, 0xD4, 0xFF)
ACCENT_RED = RGBColor(0xFF, 0x4D, 0x6D)
ACCENT_GREEN = RGBColor(0x2E, 0xD5, 0x73)
ACCENT_AMBER = RGBColor(0xFF, 0xC4, 0x4D)
CODE_BG = RGBColor(0x13, 0x1E, 0x35)
CALLOUT_BG = RGBColor(0x0F, 0x2A, 0x3F)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK_LAYOUT = prs.slide_layouts[6]


def set_bg(slide, color=BG_DARK):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=TEXT_WHITE,
             font="Helvetica Neue", align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_rich(slide, x, y, w, h, segments, *, size=18, font="Helvetica Neue",
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.3):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor

    paragraphs = [[]]
    for seg in segments:
        if seg == "NL":
            paragraphs.append([])
        else:
            paragraphs[-1].append(seg)

    for idx, segs in enumerate(paragraphs):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        for seg in segs:
            r = p.add_run()
            r.text = seg.get("text", "")
            r.font.name = seg.get("font", font)
            r.font.size = Pt(seg.get("size", size))
            r.font.bold = seg.get("bold", False)
            r.font.italic = seg.get("italic", False)
            r.font.color.rgb = seg.get("color", TEXT_WHITE)
    return tb


def add_title(slide, title_text):
    add_text(slide, Inches(0.6), Inches(0.5), Inches(12), Inches(0.9),
             title_text, size=44, bold=True, color=TEXT_WHITE)
    add_rect(slide, Inches(0.65), Inches(1.45), Inches(0.6), Inches(0.05), ACCENT_CYAN)


def add_subtitle(slide, sub_text, y=Inches(1.7)):
    add_text(slide, Inches(0.65), y, Inches(12), Inches(0.4),
             sub_text, size=16, color=TEXT_MUTED)


def add_footer(slide, footer_text):
    add_text(slide, Inches(0.65), Inches(7.0), Inches(12), Inches(0.3),
             footer_text, size=11, color=TEXT_MUTED, font="Menlo")


def add_corner_tag(slide, tag_text):
    w = Inches(2.0); h = Inches(0.45)
    x = Inches(11.0); y = Inches(0.6)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.background()
    shape.line.color.rgb = ACCENT_CYAN
    shape.line.width = Pt(1.2)
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = tag_text
    run.font.name = "Helvetica Neue"
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = ACCENT_CYAN


def add_callout(slide, x, y, w, h, text, *, accent=ACCENT_CYAN, size=20, bold=False):
    add_rect(slide, x, y, w, h, CALLOUT_BG)
    add_rect(slide, x, y, Inches(0.08), h, accent)
    tb = slide.shapes.add_textbox(x + Inches(0.35), y, w - Inches(0.5), h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Helvetica Neue"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = TEXT_WHITE


# ============================================================
# SLIDE 1 — RULER: Frontier modeli @ 256K
# ============================================================
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s)
add_title(s, "Veliki context window ≠ veliko razumijevanje")
add_subtitle(s, "RULER benchmark · frontier modeli @ 256K tokena")
add_corner_tag(s, "Research")

# Intro line
add_rich(s, Inches(0.65), Inches(2.3), Inches(12), Inches(0.6),
         [
             {"text": "Svi frontier modeli ", "size": 18, "color": TEXT_MUTED},
             {"text": "deklariraju 1M token", "size": 18, "bold": True, "color": TEXT_WHITE},
             {"text": " context window.", "size": 18, "color": TEXT_MUTED},
             "NL",
             {"text": "Na složenijim zadacima već pri ", "size": 18, "color": TEXT_MUTED},
             {"text": "256K", "size": 18, "bold": True, "color": ACCENT_CYAN},
             {"text": " (¼ deklariranog) — performans pada:", "size": 18, "color": TEXT_MUTED},
         ])

# Table
table_x = Inches(0.65)
table_y = Inches(3.6)
table_w = Inches(12.0)
row_h = Inches(0.7)

# Header row
header_y = table_y
add_rect(s, table_x, header_y, table_w, row_h, CODE_BG)
add_text(s, table_x + Inches(0.5), header_y, Inches(5.0), row_h,
         "Model", size=14, bold=True, color=TEXT_MUTED, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, table_x + Inches(6.5), header_y, Inches(2.5), row_h,
         "RULER @ 256K", size=14, bold=True, color=TEXT_MUTED,
         anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
add_text(s, table_x + Inches(9.5), header_y, Inches(2.0), row_h,
         "Status", size=14, bold=True, color=TEXT_MUTED,
         anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

# Data rows
rows = [
    ("Gemini 3 Deep Think", "84%", "drži >80%", ACCENT_GREEN),
    ("GPT-5.5", "72%", "ispod praga", ACCENT_AMBER),
    ("Claude Opus 4.7", "61%", "značajan pad", ACCENT_RED),
]
for i, (model, score, status, color) in enumerate(rows):
    y = table_y + row_h * (i + 1)
    add_rect(s, table_x, y, table_w, row_h, BG_DARK)
    # subtle divider
    add_rect(s, table_x, y, table_w, Inches(0.015), CODE_BG)
    add_text(s, table_x + Inches(0.5), y, Inches(5.0), row_h,
             model, size=18, color=TEXT_WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, table_x + Inches(6.5), y, Inches(2.5), row_h,
             score, size=22, bold=True, color=color,
             anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER, font="Menlo")
    add_text(s, table_x + Inches(9.5), y, Inches(2.0), row_h,
             status, size=14, color=color,
             anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

# Bottom callout
add_callout(s, Inches(0.65), Inches(6.7), Inches(12), Inches(0.55),
            "Deklarirano: 1M.  Realno korisno: znatno manje.",
            size=16, bold=True)

add_footer(s, "RULER · github.com/NVIDIA/RULER")


# ============================================================
# SLIDE 2 — Lost in the Middle: Definicija + natuknice
# ============================================================
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s)
add_title(s, "Lost in the Middle")
add_subtitle(s, "Stanford / UC Berkeley · 2023")
add_corner_tag(s, "Research")

# Definition box
add_callout(s, Inches(0.65), Inches(2.4), Inches(12), Inches(1.4),
            "LLM ne čita kontekst ravnomjerno. Najbolje koristi informacije "
            "na početku i kraju, a najlošije u sredini.",
            size=20, accent=ACCENT_CYAN)

# Bullets
bullets_y = Inches(4.3)
bullet_lh = Inches(0.65)

bullets = [
    ("Pozicija u kontekstu = ", "težina u pažnji", ""),
    ("Najbolja točnost na ", "rubovima", ", najlošija u sredini — U-krivulja"),
    ("Što je kontekst duži, to je ", "'mrtva zona'", " u sredini veća"),
    ("Vrijedi za sve frontier modele — ", "Claude, GPT, Gemini", ""),
]

for i, (pre, mid, post) in enumerate(bullets):
    y = bullets_y + bullet_lh * i
    add_text(s, Inches(0.65), y, Inches(0.3), Inches(0.5),
             "•", size=20, color=ACCENT_CYAN, bold=True)
    add_rich(s, Inches(1.0), y, Inches(11.5), Inches(0.6),
             [
                 {"text": pre, "size": 18, "color": TEXT_MUTED},
                 {"text": mid, "size": 18, "bold": True, "color": TEXT_WHITE},
                 {"text": post, "size": 18, "color": TEXT_MUTED},
             ])

add_footer(s, "arxiv.org/abs/2307.03172")


# ============================================================
output = "/Users/darjan.baricevic/Projects/debug-presentation/new-intro-slides.pptx"
prs.save(output)
print(f"Saved: {output}")
print(f"Slides: {len(prs.slides)}")
